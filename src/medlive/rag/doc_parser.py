
"""文档解析器：从 PDF、DOCX、PPTX、XLSX 和纯文本文件中提取文本。

这部分逻辑参考 HKUDS/LightRAG 的 document_routes.py 流程，
并改造成当前 RAG Core Service 使用的同步解析函数。
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Protocol

# ---------------------------------------------------------------------------
# docling 协议类型，避免为了类型检查强依赖 docling。
# ---------------------------------------------------------------------------


class DoclingConverter(Protocol):
    def convert(self, file_path: Path) -> DoclingResult:
        ...


class DoclingResult(Protocol):
    document: DoclingDocument


class DoclingDocument(Protocol):
    def export_to_markdown(self) -> str:
        ...


# ---------------------------------------------------------------------------
# UTF-8 文本文件扩展名。
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS: set[str] = {
    ".txt",
    ".md",
    ".mdx",
    ".html",
    ".htm",
    ".tex",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".csv",
    ".log",
    ".conf",
    ".ini",
    ".properties",
    ".sql",
    ".bat",
    ".sh",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".swift",
    ".go",
    ".rb",
    ".php",
    ".css",
    ".scss",
    ".less",
    ".rtf",  # 按文本文件兜底处理
    ".odt",  # 按文本文件兜底处理
    ".epub",  # 按文本文件兜底处理
}

# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def extract_pdf(file_bytes: bytes, password: str | None = None) -> str:
    """使用 pypdf 从 PDF 中提取文本。

    参数：
        file_bytes: PDF 原始字节。
        password: 可选解密密码。

    返回：
        提取出的文本，每页之间用换行分隔。

    异常：
        ValueError: PDF 加密且密码缺失或错误。
    """
    from pypdf import PdfReader

    pdf_file = BytesIO(file_bytes)
    reader = PdfReader(pdf_file)

    if reader.is_encrypted:
        decrypt_result = reader.decrypt(password or "")
        if decrypt_result == 0:
            if password:
                raise ValueError("Incorrect PDF password")
            raise ValueError("PDF is encrypted but no password provided")

    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def convert_pdf_docling(file_path: str | Path) -> str:
    """使用 docling 将 PDF 或其他支持的文档转换为 Markdown。

    需要额外安装 `docling`。
    """
    from docling.document_converter import DocumentConverter

    converter = DocumentConverter()
    result = converter.convert(Path(file_path))
    return result.document.export_to_markdown()


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(file_bytes: bytes) -> str:
    """按文档顺序提取 DOCX 内容，包括表格。

    表格会转换为制表符分隔的行，并用空行与段落分隔。
    """
    from docx import Document
    from docx.table import Table as DocxTable  # type: ignore[attr-defined]
    from docx.text.paragraph import Paragraph  # type: ignore[attr-defined]

    docx_file = BytesIO(file_bytes)
    doc = Document(docx_file)

    def _escape_cell(value: str | None) -> str:
        if value is None:
            return ""
        text = str(value)
        return (
            text.replace("\\", "\\\\")
            .replace("\t", "  ")
            .replace("\r\n", "<br>")
            .replace("\r", "<br>")
            .replace("\n", "<br>")
        )

    content_parts: list[str] = []
    in_table = False

    for element in doc.element.body:
        if element.tag.endswith("p"):
            if in_table:
                content_parts.append("")
                in_table = False
            paragraph = Paragraph(element, doc)
            content_parts.append(paragraph.text)
        elif element.tag.endswith("tbl"):
            if content_parts and not in_table:
                content_parts.append("")
            in_table = True
            table = DocxTable(element, doc)
            for row in table.rows:
                row_text = [_escape_cell(cell.text) for cell in row.cells]
                if any(cell for cell in row_text):
                    content_parts.append("\t".join(row_text))

    return "\n".join(content_parts)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(file_bytes: bytes) -> str:
    """提取 PPTX 所有幻灯片和形状中的文本。"""
    from pptx import Presentation

    pptx_file = BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def extract_xlsx(file_bytes: bytes) -> str:
    """提取 XLSX 内容，行内用制表符分隔，并保留工作表分隔符。"""
    from openpyxl import load_workbook

    xlsx_file = BytesIO(file_bytes)
    wb = load_workbook(xlsx_file)

    def _escape_cell(value) -> str:
        if value is None:
            return ""
        text = str(value)
        return (
            text.replace("\\", "\\\\")
            .replace("\t", "\\t")
            .replace("\r\n", "\\n")
            .replace("\r", "\\n")
            .replace("\n", "\\n")
        )

    content_parts: list[str] = []
    sep = "=" * 20

    for idx, sheet in enumerate(wb):
        if idx > 0:
            content_parts.append("")
        safe_title = str(sheet.title).replace("\n", " ").replace("\t", " ")
        content_parts.append(f"{sep} Sheet: {safe_title} {sep}")

        max_cols = sheet.max_column or 0
        for row in sheet.iter_rows(values_only=True):
            row_parts = [
                _escape_cell(row[col]) if col < len(row) else ""
                for col in range(max_cols)
            ]
            if all(p == "" for p in row_parts):
                content_parts.append("")
            else:
                content_parts.append("\t".join(row_parts))

    content_parts.append(sep)
    return "\n".join(content_parts)


# ---------------------------------------------------------------------------
# 分发器
# ---------------------------------------------------------------------------

# 二进制文件扩展名会路由到对应解析器，PDF 固定在这里处理。
BINARY_EXTENSIONS: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
}

ALL_SUPPORTED_EXTENSIONS: set[str] = TEXT_EXTENSIONS | set(BINARY_EXTENSIONS)


def parse_file_content(file_bytes: bytes, extension: str, **kwargs: object) -> str:
    """按扩展名将文件原始字节解析为文本。

    参数：
        file_bytes: 文件原始内容。
        extension: 小写扩展名，必须包含点号，例如 ``".pdf"``。
        kwargs: 传给具体解析器的额外参数，目前只有 PDF 的 ``password``。

    返回：
        提取出的文本。

    异常：
        ValueError: 扩展名不支持或解析失败。
    """
    ext = extension.lower()

    if ext in TEXT_EXTENSIONS:
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise ValueError(
                f"文件不是合法 UTF-8 文本。"
                f"请先转换为 UTF-8 后再处理: {exc}"
            ) from exc
        if not text.strip():
            raise ValueError("文件没有文本内容")
        if text.startswith("b'") or text.startswith('b"'):
            raise ValueError("文件看起来包含二进制数据")
        return text

    dispatch = BINARY_EXTENSIONS.get(ext)
    if dispatch is None:
        raise ValueError(
            f"不支持的文件扩展名: {ext}。"
            f"支持的扩展名: {sorted(ALL_SUPPORTED_EXTENSIONS)}"
        )

    if dispatch == "pdf":
        password = kwargs.get("password")
        pdf_password: str | None = str(password) if password else None
        return extract_pdf(file_bytes, password=pdf_password)
    if dispatch == "docx":
        return extract_docx(file_bytes)
    if dispatch == "pptx":
        return extract_pptx(file_bytes)
    if dispatch == "xlsx":
        return extract_xlsx(file_bytes)

    raise ValueError(f"未处理的二进制扩展名: {ext}")  # pragma: no cover
